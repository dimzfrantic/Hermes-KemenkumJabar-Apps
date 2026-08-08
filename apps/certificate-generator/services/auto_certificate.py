from __future__ import annotations

import fcntl
import json
import os
import re
import shutil
import threading
from dataclasses import dataclass
from datetime import datetime, timedelta
from pathlib import Path
from uuid import uuid4

from flask import current_app
from pptx import Presentation
from google.auth.exceptions import RefreshError, TransportError
from googleapiclient.errors import HttpError
from sqlalchemy.exc import OperationalError

from extensions import db
from models import AutoCertificateEvent, AutoCertificateItem, AutoCertificateRun
from services.certificate_photos import PHOTO_PLACEHOLDER, prepare_certificate_photo
from services.google_drive import DriveConfigurationError, upload_pdf_with_config
from services.google_sheets import ParsedSheet, SheetConfigurationError, fetch_form_rows
from services.pptx_generator import build_pdf_safe_template, build_text_replacements_from_row, convert_pptx_to_pdf_with_soffice, placeholder_token_from_header, placeholder_tokens_from_headers, replace_placeholders
from services.storage import ensure_dir, remove_dir, remove_files, slugify_filename

_SYNC_LOCKS: dict[str, threading.Lock] = {}
RETRYABLE_SYNC_EXCEPTIONS = (TransportError,)
_DB_SYNC_LOCK_PATH = Path(__file__).resolve().parent.parent / 'instance' / 'auto_certificate_sync.lock'


class AutoCertificateError(RuntimeError):
    pass



def _normalize_certificate_name(value: str | None) -> str:
    return ' '.join((value or '').split())


def _filename_timestamp(item: AutoCertificateItem) -> str:
    source_value = item.submitted_at or ''
    if not source_value and item.source_data_json:
        try:
            source_data = json.loads(item.source_data_json)
        except (TypeError, ValueError):
            source_data = {}
        if isinstance(source_data, dict):
            timestamp_header = detect_form_columns(list(source_data.keys())).get('timestamp')
            source_value = source_data.get(timestamp_header, '') if timestamp_header else ''
    safe_stamp = ''.join(ch for ch in str(source_value) if ch.isdigit())[:14]
    if safe_stamp:
        return safe_stamp
    fallback_time = item.created_at or datetime.utcnow()
    return fallback_time.strftime('%d%m%Y%H%M%S')


@dataclass
class EventSyncSummary:
    found_rows: int
    queued_rows: int
    processed_rows: int
    success_rows: int
    failed_rows: int
    message: str


def _event_lock(event_uuid: str) -> threading.Lock:
    if event_uuid not in _SYNC_LOCKS:
        _SYNC_LOCKS[event_uuid] = threading.Lock()
    return _SYNC_LOCKS[event_uuid]


class _GlobalSyncFileLock:
    def __init__(self, path: Path):
        self.path = path
        self.handle = None

    def __enter__(self):
        self.path.parent.mkdir(parents=True, exist_ok=True)
        self.handle = open(self.path, 'a+')
        try:
            fcntl.flock(self.handle.fileno(), fcntl.LOCK_EX | fcntl.LOCK_NB)
        except BlockingIOError as exc:
            self.handle.close()
            self.handle = None
            raise AutoCertificateError('Sinkronisasi global sedang berjalan. Silakan tunggu proses aktif selesai.') from exc
        self.handle.seek(0)
        self.handle.truncate()
        self.handle.write(str(os.getpid()))
        self.handle.flush()
        return self

    def __exit__(self, exc_type, exc, tb):
        if self.handle is None:
            return
        try:
            self.handle.seek(0)
            self.handle.truncate()
            fcntl.flock(self.handle.fileno(), fcntl.LOCK_UN)
        finally:
            self.handle.close()
            self.handle = None


def _normalize_header_key(value: str) -> str:
    import re

    return re.sub(r'[^a-z0-9]+', '', (value or '').strip().lower())


HEADER_ALIASES = {
    'name': ['namalengkap', 'nama', 'peserta', 'namapeserta', 'name', 'fullname', 'full name', 'participantname'],
    'email': ['email', 'alamatemail', 'surel', 'emailaddress', 'e-mail'],
    'institution': ['asalinstansi', 'instansi', 'unitkerja', 'satker', 'institution', 'organization', 'organisasi'],
    'phone': ['nomorwa', 'wa', 'nomorwa', 'nohp', 'noh p', 'telepon', 'nohpwa', 'phonenumber', 'nomortelponwa'],
    'timestamp': ['timestamp', 'stempelwaktu', 'waktu', 'tanggalwaktu'],
    'photo': ['foto', 'photo', 'pasfoto', 'uploadfoto', 'unggahfoto', 'filefoto', 'linkfoto', 'fotopeserta'],
}


def detect_form_columns(headers: list[str]) -> dict[str, str | None]:
    normalized_map = {_normalize_header_key(header): header for header in headers}
    result: dict[str, str | None] = {}
    for key, aliases in HEADER_ALIASES.items():
        result[key] = next((normalized_map[alias] for alias in aliases if alias in normalized_map), None)
    return result


def _extract_template_placeholders(template_path: str) -> list[str]:
    prs = Presentation(template_path)
    placeholders: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, 'has_text_frame', False):
                continue
            for match in re.findall(r'\{\{[^{}]+\}\}', shape.text or ''):
                if match not in placeholders:
                    placeholders.append(match)
    return placeholders


def _infer_columns_from_template(headers: list[str], template_path: str, columns: dict[str, str | None]) -> dict[str, str | None]:
    template_tokens = set(_extract_template_placeholders(template_path))
    template_header_matches = [header for header in headers if placeholder_token_from_header(header) in template_tokens]
    reserved = {value for value in columns.values() if value}

    def pick_candidate(preferred_keys: list[str] | None = None) -> str | None:
        preferred_keys = preferred_keys or []
        normalized_lookup = {_normalize_header_key(header): header for header in template_header_matches if header not in reserved}
        for key in preferred_keys:
            candidate = normalized_lookup.get(_normalize_header_key(key))
            if candidate:
                return candidate
        for header in template_header_matches:
            if header not in reserved:
                return header
        return None

    if not columns.get('name'):
        columns['name'] = pick_candidate(['name', 'fullname', 'full name', 'nama lengkap', 'nama', 'participant name'])
        if columns.get('name'):
            reserved.add(columns['name'])
    if not columns.get('institution'):
        columns['institution'] = pick_candidate(['institution', 'organization', 'organisasi', 'instansi', 'asal instansi', 'unit kerja'])
        if columns.get('institution'):
            reserved.add(columns['institution'])
    if not columns.get('email'):
        columns['email'] = pick_candidate(['email', 'email address', 'alamat email', 'e-mail'])
        if columns.get('email'):
            reserved.add(columns['email'])
    if not columns.get('phone'):
        columns['phone'] = pick_candidate(['nomor telpon / wa', 'no wa', 'wa', 'phone', 'phone number'])
        if columns.get('phone'):
            reserved.add(columns['phone'])
    if not columns.get('timestamp'):
        columns['timestamp'] = pick_candidate(['timestamp', 'stempel waktu'])
        if columns.get('timestamp'):
            reserved.add(columns['timestamp'])
    return columns


def validate_event_configuration(parsed: ParsedSheet, template_path: str) -> dict:
    columns = detect_form_columns(parsed.headers)
    columns = _infer_columns_from_template(parsed.headers, template_path, columns)
    if not columns.get('name'):
        raise AutoCertificateError('Kolom nama peserta tidak terdeteksi dari Google Form Response maupun placeholder template. Gunakan placeholder yang sama persis dengan header sheet, misalnya {{Name}} atau {{Nama Lengkap}}.')

    preview_root = ensure_dir(Path(template_path).parent / '_automation_validation')
    safe_template_path = preview_root / 'template-safe-check.pptx'
    build_pdf_safe_template(
        template_path,
        str(safe_template_path),
        placeholders=placeholder_tokens_from_headers(parsed.headers) + [PHOTO_PLACEHOLDER],
        working_dir=str(preview_root / 'build'),
        soffice_path=current_app.config.get('SOFFICE_PATH', ''),
    )
    sample_row = next((row for row in parsed.rows if (row.get(columns['name'] or '') or '').strip()), None)
    return {
        'columns': columns,
        'sample_name': sample_row.get(columns['name'] or '', '') if sample_row else '',
        'total_rows': len(parsed.rows),
        'has_sample_row': sample_row is not None,
    }


def _participant_key(event: AutoCertificateEvent, row: dict) -> str:
    ts_value = (row.get(event.timestamp_column or '') or '').strip()
    email_value = (row.get(event.email_column or '') or '').strip().lower()
    if ts_value and email_value:
        return f'{ts_value}|{email_value}'
    if ts_value:
        return f'{ts_value}|row-{row.get("_row_number", 0)}'
    if email_value:
        return f'{email_value}|row-{row.get("_row_number", 0)}'
    return f'row-{row.get("_row_number", 0)}|{(row.get(event.name_column or "") or "").strip().lower()}'


def _recalculate_event_counters(event: AutoCertificateEvent):
    items = AutoCertificateItem.query.filter_by(event_id=event.id).all()
    event.total_responses = len(items)
    event.pending_responses = sum(1 for item in items if item.status == 'pending')
    event.processing_responses = sum(1 for item in items if item.status == 'processing')
    event.success_responses = sum(1 for item in items if item.status == 'success')
    event.failed_responses = sum(1 for item in items if item.status == 'failed')


def _processing_timeout_minutes(event: AutoCertificateEvent) -> int:
    interval = max(int(event.polling_interval_minutes or 5), 1)
    return max(interval * 2, 10)


def _recover_stuck_processing_items(event: AutoCertificateEvent) -> int:
    cutoff = datetime.utcnow() - timedelta(minutes=_processing_timeout_minutes(event))
    stuck_items = (
        AutoCertificateItem.query
        .filter_by(event_id=event.id, status='processing')
        .filter(AutoCertificateItem.processed_at.is_(None))
        .filter(AutoCertificateItem.drive_file_id.is_(None))
        .filter(AutoCertificateItem.updated_at <= cutoff)
        .all()
    )
    recovered = 0
    for item in stuck_items:
        item.status = 'pending'
        item.error_message = 'Antrean dipulihkan otomatis setelah proses sebelumnya terputus.'
        item.processed_at = None
        recovered += 1
    if recovered:
        _recalculate_event_counters(event)
        db.session.commit()
    return recovered


def _mark_item_for_retry(item: AutoCertificateItem, message: str | None = None):
    item.status = 'pending'
    item.error_message = message or 'Akan dicoba ulang otomatis pada sinkronisasi berikutnya.'
    item.processed_at = None
    item.output_filename = None
    item.drive_file_id = None
    item.drive_link = None


def _retryable_http_error(exc: HttpError) -> bool:
    status_code = getattr(getattr(exc, 'resp', None), 'status', None)
    return status_code in {408, 409, 425, 429, 500, 502, 503, 504}


def _is_retryable_generation_error(exc: Exception) -> bool:
    if isinstance(exc, RETRYABLE_SYNC_EXCEPTIONS):
        return True
    if isinstance(exc, HttpError):
        return _retryable_http_error(exc)
    return False


def _prepare_event_template(event: AutoCertificateEvent, headers: list[str]) -> Path:
    event_root = ensure_dir(Path(current_app.config['AUTO_EVENT_DIR']) / event.event_uuid)
    template_original = event_root / 'template-original.pptx'
    template_safe = event_root / 'template-safe.pptx'
    if not template_original.exists():
        raise AutoCertificateError('Template kegiatan tidak ditemukan pada server.')
    if event.mapping_json is not None:
        try:
            mapping = json.loads(event.mapping_json or '{}')
        except (TypeError, ValueError):
            mapping = {}
        placeholder_tokens = list(dict.fromkeys(target for target in mapping.values() if target))
        if not placeholder_tokens:
            shutil.copy2(template_original, template_safe)
            return template_safe
        build_pdf_safe_template(
            str(template_original),
            str(template_safe),
            placeholders=placeholder_tokens,
            working_dir=str(event_root / '_safe_build'),
            soffice_path=current_app.config.get('SOFFICE_PATH', ''),
            cleanup_working_dir=True,
        )
        return template_safe

    # Kompatibilitas kegiatan lama yang belum mempunyai mapping dinamis.
    legacy_replacements = {}
    if event.name_column:
        legacy_replacements['{{nama}}'] = f'{{{{{event.name_column}}}}}'
    if event.institution_column:
        legacy_replacements['{{instansi}}'] = f'{{{{{event.institution_column}}}}}'
    template_source = template_original
    migrated_template = event_root / 'template-original-migrated.pptx'
    if legacy_replacements:
        replace_placeholders(str(template_original), str(migrated_template), legacy_replacements)
        template_source = migrated_template
    build_pdf_safe_template(
        str(template_source),
        str(template_safe),
        placeholders=placeholder_tokens_from_headers(headers) + [PHOTO_PLACEHOLDER],
        working_dir=str(event_root / '_safe_build'),
        soffice_path=current_app.config.get('SOFFICE_PATH', ''),
        cleanup_working_dir=True,
    )
    return template_safe


def _generate_item_certificate(event: AutoCertificateEvent, item: AutoCertificateItem, template_safe_path: str) -> dict:
    runtime_root = ensure_dir(Path(current_app.config['RUNTIME_DIR']) / 'auto-events' / event.event_uuid)
    normalized_name = _normalize_certificate_name(item.participant_name)
    base_name = slugify_filename(f'Sertifikat {normalized_name}', default=f'Sertifikat Peserta {item.source_row_number}')
    safe_stamp = _filename_timestamp(item)
    base_name = slugify_filename(f'{base_name} - {safe_stamp}', default=base_name)
    pptx_path = runtime_root / f'{uuid4().hex}-{base_name}.pptx'
    pdf_path = runtime_root / f'{pptx_path.stem}.pdf'
    photo_path = None
    try:
        try:
            source_data = json.loads(item.source_data_json or '{}')
        except (TypeError, ValueError):
            source_data = {}
        if not isinstance(source_data, dict):
            source_data = {}
        if not source_data:
            fallback_values = {
                event.name_column: item.participant_name,
                event.institution_column: item.institution_name,
                event.email_column: item.email,
                event.phone_column: item.phone,
                event.photo_column: item.photo_url,
                event.timestamp_column: item.submitted_at,
            }
            source_data = {key: value for key, value in fallback_values.items() if key}
        try:
            mapping = json.loads(event.mapping_json or '{}') if event.mapping_json is not None else None
        except (TypeError, ValueError):
            mapping = {}
        if mapping is None:
            text_replacements = build_text_replacements_from_row(source_data)
        else:
            text_replacements = {
                target: ('' if source_data.get(source) is None else str(source_data.get(source)).strip())
                for source, target in mapping.items()
                if target and target != PHOTO_PLACEHOLDER
            }
        image_replacements = {}
        photo_path = prepare_certificate_photo(
            item.photo_url,
            runtime_root,
            google_token_path=current_app.config['GOOGLE_TOKEN_PATH'],
            drive_scopes=list(current_app.config['DRIVE_SCOPES']),
            row_number=item.source_row_number,
            use_cached_service=True,
        ) if item.photo_url else None
        if photo_path:
            image_replacements[PHOTO_PLACEHOLDER] = photo_path
        replace_placeholders(
            template_safe_path,
            str(pptx_path),
            text_replacements,
            image_replacements=image_replacements,
        )
        actual_pdf = convert_pptx_to_pdf_with_soffice(str(pptx_path), str(runtime_root), current_app.config.get('SOFFICE_PATH', ''))
        filename = f'{base_name}.pdf'
        upload_result = upload_pdf_with_config(
            actual_pdf,
            event.drive_folder_id,
            filename,
            current_app.config['GOOGLE_TOKEN_PATH'],
            list(current_app.config['DRIVE_SCOPES']),
            use_cached_service=True,
        )
        return {
            'ok': True,
            'filename': filename,
            'drive_file_id': upload_result.get('id'),
            'drive_link': upload_result.get('webViewLink'),
            'message': 'Upload berhasil.',
        }
    finally:
        remove_files([pptx_path, pdf_path, photo_path])
        remove_dir(runtime_root / '_soffice_profiles')


def sync_event(event_id: int) -> EventSyncSummary:
    event = db.session.get(AutoCertificateEvent, event_id)
    if event is None:
        raise AutoCertificateError('Kegiatan automation tidak ditemukan.')
    lock = _event_lock(event.event_uuid)
    if not lock.acquire(blocking=False):
        raise AutoCertificateError('Sinkronisasi kegiatan sedang berjalan. Silakan tunggu proses yang aktif selesai.')
    run = None
    try:
        with _GlobalSyncFileLock(_DB_SYNC_LOCK_PATH):
            event = db.session.get(AutoCertificateEvent, event_id)
            if event is None:
                raise AutoCertificateError('Kegiatan automation tidak ditemukan.')

            stale_runs = AutoCertificateRun.query.filter_by(event_id=event.id, status='running').all()
            for stale_run in stale_runs:
                stale_run.finished_at = datetime.utcnow()
                stale_run.status = 'aborted'
                stale_run.message = 'Run sebelumnya ditutup otomatis karena ada proses sinkronisasi baru yang mengambil alih.'
            if stale_runs:
                db.session.commit()

            event.last_run_started_at = datetime.utcnow()
            event.status = 'active' if event.enabled else event.status
            run = AutoCertificateRun(event_id=event.id, status='running')
            db.session.add(run)
            db.session.commit()

            recovered_rows = _recover_stuck_processing_items(event)

            parsed = fetch_form_rows(event.spreadsheet_id, worksheet_name=event.worksheet_name)
            event.spreadsheet_title = parsed.spreadsheet_title
            event.worksheet_name = parsed.selected_sheet
            detected_columns = detect_form_columns(parsed.headers)
            if event.mapping_json is None:
                detected_columns = _infer_columns_from_template(
                    parsed.headers,
                    str(Path(current_app.config['AUTO_EVENT_DIR']) / event.event_uuid / 'template-original.pptx'),
                    detected_columns,
                )
            if event.mapping_json is not None and event.name_column not in parsed.headers:
                raise AutoCertificateError(
                    f'Kolom sumber nama file "{event.name_column}" tidak ditemukan pada Google Sheet.'
                )
            if event.mapping_json is None and (not event.name_column or event.name_column not in parsed.headers):
                event.name_column = detected_columns.get('name')
            if event.mapping_json is None and (not event.institution_column or event.institution_column not in parsed.headers):
                event.institution_column = detected_columns.get('institution')
            if event.mapping_json is None and (not event.email_column or event.email_column not in parsed.headers):
                event.email_column = detected_columns.get('email')
            if event.mapping_json is None and (not event.phone_column or event.phone_column not in parsed.headers):
                event.phone_column = detected_columns.get('phone')
            if event.mapping_json is None and (not event.timestamp_column or event.timestamp_column not in parsed.headers):
                event.timestamp_column = detected_columns.get('timestamp')
            if event.mapping_json is None and (not event.photo_column or event.photo_column not in parsed.headers):
                event.photo_column = detected_columns.get('photo')
            if not event.name_column:
                raise AutoCertificateError('Kolom nama peserta pada auto generate tidak terdeteksi dari sheet maupun placeholder template.')
            run.found_rows = len(parsed.rows)

            existing = {item.participant_key: item for item in AutoCertificateItem.query.filter_by(event_id=event.id).all()}
            queued_rows = 0
            for row in parsed.rows:
                key = _participant_key(event, row)
                source_data = {str(header): value for header, value in row.items() if not str(header).startswith('_')}
                source_data_json = json.dumps(source_data, ensure_ascii=False)
                if key in existing:
                    item = existing[key]
                    item.source_data_json = source_data_json
                    item.submitted_at = (row.get(event.timestamp_column or '') or '').strip() or None
                    item.participant_name = (row.get(event.name_column or '') or '').strip() or None
                    item.institution_name = (row.get(event.institution_column or '') or '').strip() or None
                    item.email = (row.get(event.email_column or '') or '').strip() or None
                    item.phone = (row.get(event.phone_column or '') or '').strip() or None
                    if event.photo_column:
                        item.photo_url = (row.get(event.photo_column or '') or '').strip() or None
                    continue
                item = AutoCertificateItem(
                    event_id=event.id,
                    participant_key=key,
                    source_row_number=row.get('_row_number', 0),
                    submitted_at=(row.get(event.timestamp_column or '') or '').strip() or None,
                    participant_name=(row.get(event.name_column or '') or '').strip() or None,
                    institution_name=(row.get(event.institution_column or '') or '').strip() or None,
                    email=(row.get(event.email_column or '') or '').strip() or None,
                    phone=(row.get(event.phone_column or '') or '').strip() or None,
                    photo_url=(row.get(event.photo_column or '') or '').strip() or None,
                    source_data_json=source_data_json,
                    status='pending',
                )
                db.session.add(item)
                queued_rows += 1
            db.session.commit()
            run.queued_rows = queued_rows

            event_root = ensure_dir(Path(current_app.config['AUTO_EVENT_DIR']) / event.event_uuid)
            template_safe_path = _prepare_event_template(event, parsed.headers)

            pending_items = AutoCertificateItem.query.filter_by(event_id=event.id, status='pending').order_by(AutoCertificateItem.id.asc()).all()

            success_rows = 0
            failed_rows = 0
            for item in pending_items:
                try:
                    item.status = 'processing'
                    item.error_message = None
                    item.processed_at = None
                    db.session.commit()

                    result = _generate_item_certificate(event, item, str(template_safe_path))
                    item.status = 'success'
                    item.output_filename = result.get('filename')
                    item.drive_file_id = result.get('drive_file_id')
                    item.drive_link = result.get('drive_link')
                    item.error_message = None
                    item.processed_at = datetime.utcnow()
                    success_rows += 1
                except Exception as exc:
                    if _is_retryable_generation_error(exc):
                        _mark_item_for_retry(item, 'Koneksi internet/Google sempat terputus. Peserta dikembalikan ke antrean untuk dicoba ulang otomatis.')
                        db.session.commit()
                        raise AutoCertificateError('Koneksi internet/Google terputus di tengah proses. Sisa antrean akan dicoba ulang otomatis pada sinkronisasi berikutnya.') from exc
                    item.status = 'failed'
                    item.error_message = str(exc)
                    item.processed_at = datetime.utcnow()
                    failed_rows += 1
                db.session.commit()

            _recalculate_event_counters(event)
            event.last_synced_at = datetime.utcnow()
            event.last_run_finished_at = datetime.utcnow()
            event.next_run_at = datetime.utcnow() + timedelta(minutes=event.polling_interval_minutes or 5)
            event.last_error = None
            event.status = 'active' if event.enabled else 'inactive'

            run.finished_at = datetime.utcnow()
            run.status = 'completed' if failed_rows == 0 else 'completed_with_errors'
            run.processed_rows = len(pending_items)
            run.success_rows = success_rows
            run.failed_rows = failed_rows
            recovery_note = f', dipulihkan: {recovered_rows}' if recovered_rows else ''
            run.message = f'Sinkron selesai. Baru: {queued_rows}, sukses: {success_rows}, gagal: {failed_rows}{recovery_note}.'
            db.session.commit()
            return EventSyncSummary(
                found_rows=len(parsed.rows),
                queued_rows=queued_rows,
                processed_rows=len(pending_items),
                success_rows=success_rows,
                failed_rows=failed_rows,
                message=run.message,
            )
    except OperationalError as exc:
        db.session.rollback()
        if 'database is locked' in str(exc).lower():
            raise AutoCertificateError('Database sedang dipakai proses sinkronisasi lain. Silakan tunggu sebentar lalu coba lagi.') from exc
        raise AutoCertificateError(str(exc)) from exc
    except (SheetConfigurationError, DriveConfigurationError, RefreshError, HttpError, AutoCertificateError, RuntimeError, TransportError) as exc:
        db.session.rollback()
        event = db.session.get(AutoCertificateEvent, event_id)
        if event is not None:
            event.status = 'error'
            if isinstance(exc, TransportError):
                event.last_error = 'Server tidak dapat terhubung ke layanan Google saat proses sinkronisasi.'
            else:
                event.last_error = str(exc)
            event.last_run_finished_at = datetime.utcnow()
            event.next_run_at = datetime.utcnow() + timedelta(minutes=event.polling_interval_minutes or 5)
            _recalculate_event_counters(event)
        if run is None:
            run = AutoCertificateRun.query.filter_by(event_id=event_id).order_by(AutoCertificateRun.id.desc()).first()
        if run and run.status == 'running':
            run.finished_at = datetime.utcnow()
            run.status = 'failed'
            run.message = str(exc)
        db.session.commit()
        raise AutoCertificateError(str(exc)) from exc
    finally:
        lock.release()


def retry_failed_items(event_id: int) -> int:
    event = db.session.get(AutoCertificateEvent, event_id)
    if event is None:
        raise AutoCertificateError('Kegiatan automation tidak ditemukan.')
    items = AutoCertificateItem.query.filter_by(event_id=event.id, status='failed').all()
    count = 0
    for item in items:
        item.status = 'pending'
        item.error_message = None
        item.processed_at = None
        count += 1
    _recalculate_event_counters(event)
    db.session.commit()
    return count


def reset_processing_items(event_id: int, *, only_stale: bool = False) -> int:
    event = db.session.get(AutoCertificateEvent, event_id)
    if event is None:
        raise AutoCertificateError('Kegiatan automation tidak ditemukan.')
    query = AutoCertificateItem.query.filter_by(event_id=event.id, status='processing')
    if only_stale:
        cutoff = datetime.utcnow() - timedelta(minutes=_processing_timeout_minutes(event))
        query = (
            query
            .filter(AutoCertificateItem.processed_at.is_(None))
            .filter(AutoCertificateItem.drive_file_id.is_(None))
            .filter(AutoCertificateItem.updated_at <= cutoff)
        )
    items = query.all()
    count = 0
    for item in items:
        _mark_item_for_retry(item, 'Antrean direset manual untuk diproses ulang.')
        count += 1
    _recalculate_event_counters(event)
    db.session.commit()
    return count
