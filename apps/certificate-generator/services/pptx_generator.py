from __future__ import annotations

import shutil
import subprocess
from pathlib import Path
from uuid import uuid4

from flask import current_app
from PIL import Image
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu


def placeholder_token_from_header(header: str) -> str:
    text = '' if header is None else str(header).strip()
    if not text or text.startswith('_'):
        return ''
    return f'{{{{{text}}}}}'


def placeholder_tokens_from_headers(headers: list[str] | None) -> list[str]:
    tokens: list[str] = []
    for header in headers or []:
        token = placeholder_token_from_header(header)
        if token and token not in tokens:
            tokens.append(token)
    return tokens


def build_text_replacements_from_row(row: dict) -> dict[str, str]:
    replacements: dict[str, str] = {}
    for header, value in (row or {}).items():
        token = placeholder_token_from_header(header)
        if not token:
            continue
        replacements[token] = '' if value is None else str(value).strip()
    return replacements


def expand_placeholder_tokens(placeholders: list[str] | None) -> list[str]:
    expanded: list[str] = []
    for token in placeholders or []:
        if token not in expanded:
            expanded.append(token)
    return expanded


def _normalize_replacements(replacements: dict[str, str]) -> dict[str, str]:
    normalized = {}
    for key, value in replacements.items():
        text = '' if value is None else str(value).strip()
        normalized[key] = text
    return normalized


def _replace_in_runs(paragraph, replacements: dict[str, str]) -> bool:
    if not paragraph.runs:
        return False

    changed = False
    for old, new in replacements.items():
        while True:
            runs = list(paragraph.runs)
            full_text = ''.join(run.text for run in runs)
            start = full_text.find(old)
            if start == -1:
                break

            end = start + len(old)
            run_ranges = []
            cursor = 0
            for idx, run in enumerate(runs):
                run_text = run.text or ''
                next_cursor = cursor + len(run_text)
                run_ranges.append((idx, cursor, next_cursor))
                cursor = next_cursor

            impacted = []
            for idx, run_start, run_end in run_ranges:
                overlap_start = max(start, run_start)
                overlap_end = min(end, run_end)
                if overlap_start < overlap_end:
                    impacted.append((idx, run_start, run_end, overlap_end - overlap_start))

            if not impacted:
                break

            anchor_idx = max(impacted, key=lambda item: item[3])[0]
            for idx, run_start, run_end, _ in impacted:
                run = runs[idx]
                original = run.text or ''
                local_token_start = max(start, run_start) - run_start
                local_token_end = min(end, run_end) - run_start
                prefix = original[:local_token_start]
                suffix = original[local_token_end:]
                run.text = f'{prefix}{new}{suffix}' if idx == anchor_idx else f'{prefix}{suffix}'
            changed = True
    return changed


def _replace_paragraph_text_preserve_style(paragraph, new_text: str):
    runs = list(paragraph.runs)
    if not runs:
        paragraph.text = new_text
        return

    anchor_run = runs[0]
    anchor_font = anchor_run.font
    font_name = anchor_font.name
    font_size = anchor_font.size
    font_bold = anchor_font.bold
    font_italic = anchor_font.italic
    font_underline = anchor_font.underline
    color_rgb = None
    try:
        if anchor_font.color.type is not None and anchor_font.color.rgb is not None:
            color_rgb = anchor_font.color.rgb
    except Exception:
        pass

    anchor_run.text = new_text
    for extra_run in runs[1:]:
        extra_run.text = ''

    anchor_font = anchor_run.font
    anchor_font.name = font_name
    anchor_font.size = font_size
    anchor_font.bold = font_bold
    anchor_font.italic = font_italic
    anchor_font.underline = font_underline
    try:
        if color_rgb is not None:
            anchor_font.color.rgb = color_rgb
    except Exception:
        pass



def _replace_in_shape(shape, replacements: dict[str, str]):
    for paragraph in shape.text_frame.paragraphs:
        changed = _replace_in_runs(paragraph, replacements)
        if not changed:
            raw_text = paragraph.text or ''
            updated = raw_text
            for old, new in replacements.items():
                updated = updated.replace(old, new)
            if updated != raw_text:
                _replace_paragraph_text_preserve_style(paragraph, updated)


def _delete_shape(shape):
    element = shape._element
    element.getparent().remove(element)


def _apply_center_crop(picture, image_path: str, box_width, box_height):
    try:
        with Image.open(image_path) as img:
            image_width, image_height = img.size
    except Exception:
        return
    if not image_width or not image_height or not box_width or not box_height:
        return
    image_ratio = image_width / image_height
    box_ratio = int(box_width) / int(box_height)
    if image_ratio > box_ratio:
        crop = max(0, min(0.49, (1 - (box_ratio / image_ratio)) / 2))
        picture.crop_left = crop
        picture.crop_right = crop
    elif image_ratio < box_ratio:
        crop = max(0, min(0.49, (1 - (image_ratio / box_ratio)) / 2))
        picture.crop_top = crop
        picture.crop_bottom = crop


def _resolved_image_box(token: str, left, top, width, height):
    if token != '{{foto}}':
        return left, top, width, height
    if not width or not height:
        return left, top, width, height
    try:
        box_ratio = int(width) / int(height)
    except Exception:
        return left, top, width, height
    if box_ratio <= 1.2:
        return left, top, width, height
    portrait_height = int(int(width) * 4 / 3)
    if portrait_height <= int(height):
        return left, top, width, height
    center_y = int(top) + (int(height) // 2)
    adjusted_top = center_y - (portrait_height // 2)
    return left, adjusted_top, width, portrait_height


def _boxes_overlap(left_a, top_a, width_a, height_a, left_b, top_b, width_b, height_b) -> bool:
    right_a = int(left_a) + int(width_a)
    bottom_a = int(top_a) + int(height_a)
    right_b = int(left_b) + int(width_b)
    bottom_b = int(top_b) + int(height_b)
    return not (
        right_a <= int(left_b)
        or right_b <= int(left_a)
        or bottom_a <= int(top_b)
        or bottom_b <= int(top_a)
    )


def _collect_overlapping_placeholder_shapes(slide, current_shape, token: str, left, top, width, height):
    overlaps = []
    if token != '{{foto}}':
        return overlaps
    for other_shape in slide.shapes:
        if other_shape == current_shape:
            continue
        if not getattr(other_shape, 'has_text_frame', False):
            continue
        other_text = (other_shape.text or '').strip()
        if '{{' not in other_text or other_text == token:
            continue
        if _boxes_overlap(left, top, width, height, other_shape.left, other_shape.top, other_shape.width, other_shape.height):
            overlaps.append(other_shape)
    overlaps.sort(key=lambda shape: int(shape.top))
    return overlaps


def _shrink_box_to_fit_text_below(left, top, width, height, overlapping_shapes, slide_height):
    if not overlapping_shapes or slide_height is None:
        return left, top, width, height
    gap = max(91440 // 4, int(height) // 12)
    text_stack_height = sum(int(shape.height) for shape in overlapping_shapes)
    text_stack_height += gap * len(overlapping_shapes)
    available_height = int(slide_height) - int(top) - text_stack_height
    if available_height <= 0 or available_height >= int(height):
        return left, top, width, height
    scale = available_height / max(int(height), 1)
    center_x = int(left) + (int(width) // 2)
    new_width = max(1, int(int(width) * scale))
    new_height = max(1, int(available_height))
    new_left = center_x - (new_width // 2)
    return new_left, int(top), new_width, new_height


def _stack_overlapping_placeholder_shapes_below(photo_bottom: int, overlapping_shapes):
    if not overlapping_shapes:
        return
    tallest = max(int(shape.height) for shape in overlapping_shapes)
    gap = max(91440 // 4, tallest // 12)
    cursor = int(photo_bottom) + gap
    for shape in overlapping_shapes:
        shape.top = Emu(cursor)
        cursor += int(shape.height) + gap


def _replace_image_placeholders(prs: Presentation, image_replacements: dict[str, str]):
    normalized = {token: str(path) for token, path in (image_replacements or {}).items() if path and Path(str(path)).exists()}
    if not normalized:
        return
    for slide in prs.slides:
        for shape in list(slide.shapes):
            if not getattr(shape, 'has_text_frame', False):
                continue
            text = shape.text or ''
            matched_token = next((token for token in normalized if token in text), None)
            if not matched_token:
                continue
            image_path = normalized[matched_token]
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            left, top, width, height = _resolved_image_box(matched_token, left, top, width, height)
            overlapping_shapes = _collect_overlapping_placeholder_shapes(slide, shape, matched_token, left, top, width, height)
            left, top, width, height = _shrink_box_to_fit_text_below(left, top, width, height, overlapping_shapes, prs.slide_height)
            overlapping_shapes = _collect_overlapping_placeholder_shapes(slide, shape, matched_token, left, top, width, height)
            _stack_overlapping_placeholder_shapes_below(int(top) + int(height), overlapping_shapes)
            _delete_shape(shape)
            picture = slide.shapes.add_picture(image_path, left, top, width=width, height=height)
            _apply_center_crop(picture, image_path, width, height)


def replace_placeholders(
    input_pptx: str,
    output_pptx: str,
    replacements: dict[str, str],
    image_replacements: dict[str, str] | None = None,
):
    prs = Presentation(input_pptx)
    _replace_image_placeholders(prs, image_replacements or {})
    normalized = _normalize_replacements(replacements)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, 'has_text_frame', False):
                continue
            _replace_in_shape(shape, normalized)
    prs.save(output_pptx)


def _find_placeholder_shapes(prs: Presentation, placeholders: list[str]) -> list[tuple[int, int]]:
    found: list[tuple[int, int]] = []
    for slide_index, slide in enumerate(prs.slides):
        for shape_index, shape in enumerate(slide.shapes):
            if not getattr(shape, 'has_text_frame', False):
                continue
            text = shape.text or ''
            if any(token in text for token in placeholders):
                found.append((slide_index, shape_index))
    return found


def _validate_placeholder_shapes(prs: Presentation, placeholders: list[str]) -> list[tuple[int, int]]:
    if len(prs.slides) != 1:
        raise RuntimeError('Template sertifikat aman PDF hanya mendukung 1 slide. Gabungkan desain menjadi satu slide dengan background PNG dan textbox placeholder.')

    found = _find_placeholder_shapes(prs, placeholders)
    if not found:
        raise RuntimeError(
            'Placeholder template tidak ditemukan. Pastikan template PPTX memuat placeholder yang sama persis dengan header sumber data, misalnya {{Nama Lengkap}}.'
        )

    for slide_index, shape_index in found:
        shape = prs.slides[slide_index].shapes[shape_index]
        if not getattr(shape, 'has_text_frame', False):
            raise RuntimeError('Placeholder harus berada pada textbox teks, bukan pada shape/grafik lain.')
        text = shape.text or ''
        cleaned = text
        for token in placeholders:
            cleaned = cleaned.replace(token, '')
        if cleaned.strip():
            raise RuntimeError(
                'Textbox placeholder hanya boleh berisi satu placeholder murni per textbox. '
                'Gunakan satu textbox terpisah untuk setiap placeholder header sumber data; {{foto}} tetap khusus untuk gambar. '
                'Pindahkan ornamen/desain ke background PNG dan sisakan textbox untuk placeholder saja.'
            )
    return found


_XML_PARAGRAPH_ALIGNMENTS = {
    'l': PP_ALIGN.LEFT,
    'ctr': PP_ALIGN.CENTER,
    'r': PP_ALIGN.RIGHT,
    'just': PP_ALIGN.JUSTIFY,
    'dist': PP_ALIGN.DISTRIBUTE,
    'thaiDist': PP_ALIGN.THAI_DISTRIBUTE,
    'justLow': PP_ALIGN.JUSTIFY_LOW,
}


def _list_style_alignment(shape, paragraph_level: int):
    try:
        level_nodes = shape._element.txBody.xpath(f'./a:lstStyle/a:lvl{paragraph_level + 1}pPr')
        if level_nodes:
            return _XML_PARAGRAPH_ALIGNMENTS.get(level_nodes[0].get('algn'))
    except Exception:
        return None
    return None


def _effective_paragraph_alignment(src_shape, src_paragraph):
    if src_paragraph.alignment is not None:
        return src_paragraph.alignment

    alignment = _list_style_alignment(src_shape, src_paragraph.level)
    if alignment is not None:
        return alignment

    if getattr(src_shape, 'is_placeholder', False):
        try:
            placeholder_idx = src_shape.placeholder_format.idx
            slide_layout = src_shape.part.slide.slide_layout
            for layout_shape in slide_layout.placeholders:
                if layout_shape.placeholder_format.idx != placeholder_idx:
                    continue
                alignment = _list_style_alignment(layout_shape, src_paragraph.level)
                if alignment is not None:
                    return alignment
                layout_paragraphs = list(layout_shape.text_frame.paragraphs)
                if layout_paragraphs and layout_paragraphs[0].alignment is not None:
                    return layout_paragraphs[0].alignment
        except Exception:
            pass
    return None


def _copy_textbox_style(src_shape, dest_shape):
    dest_shape.text_frame.clear()
    dest_shape.text_frame.word_wrap = src_shape.text_frame.word_wrap
    dest_shape.text_frame.auto_size = src_shape.text_frame.auto_size
    dest_shape.text_frame.margin_left = src_shape.text_frame.margin_left
    dest_shape.text_frame.margin_right = src_shape.text_frame.margin_right
    dest_shape.text_frame.margin_top = src_shape.text_frame.margin_top
    dest_shape.text_frame.margin_bottom = src_shape.text_frame.margin_bottom
    dest_shape.text_frame.vertical_anchor = src_shape.text_frame.vertical_anchor

    dest_shape.fill.background()
    dest_shape.line.fill.background()

    src_paragraphs = list(src_shape.text_frame.paragraphs)
    for paragraph_index, src_paragraph in enumerate(src_paragraphs):
        dest_paragraph = dest_shape.text_frame.paragraphs[0] if paragraph_index == 0 else dest_shape.text_frame.add_paragraph()
        dest_paragraph.alignment = _effective_paragraph_alignment(src_shape, src_paragraph)
        dest_paragraph.level = src_paragraph.level
        try:
            dest_paragraph.line_spacing = src_paragraph.line_spacing
        except Exception:
            pass
        try:
            dest_paragraph.space_before = src_paragraph.space_before
        except Exception:
            pass
        try:
            dest_paragraph.space_after = src_paragraph.space_after
        except Exception:
            pass

        if not src_paragraph.runs:
            dest_paragraph.text = src_paragraph.text
            continue

        for src_run in src_paragraph.runs:
            dest_run = dest_paragraph.add_run()
            dest_run.text = src_run.text
            src_font = src_run.font
            dest_font = dest_run.font
            dest_font.name = src_font.name
            dest_font.size = src_font.size
            dest_font.bold = src_font.bold
            dest_font.italic = src_font.italic
            dest_font.underline = src_font.underline
            try:
                if src_font.color.type is not None and src_font.color.rgb is not None:
                    dest_font.color.rgb = src_font.color.rgb
            except Exception:
                pass


def resolve_soffice_path(configured: str = '') -> str:
    if configured:
        return configured
    for candidate in ('soffice', '/usr/bin/soffice', '/snap/bin/libreoffice'):
        if shutil.which(candidate) or Path(candidate).exists():
            return candidate
    raise RuntimeError('LibreOffice/soffice tidak ditemukan pada server.')


def _resolve_soffice() -> str:
    return resolve_soffice_path(current_app.config.get('SOFFICE_PATH', ''))


def _compress_png_for_pdf_target(
    png_path: str,
    max_bytes: int = 200 * 1024,
    min_width: int = 1800,
    min_height: int = 1200,
) -> str:
    path = Path(png_path)
    if not path.exists():
        return png_path

    with Image.open(path) as img:
        image = img.convert('RGB')
        width, height = image.size
        quality_steps = [92, 88, 84, 80, 76, 72]
        scale_steps = [1.0, 0.9, 0.82, 0.75, 0.68]
        best_candidate = None

        for scale in scale_steps:
            scaled_width = max(min_width, int(width * scale))
            scaled_height = max(min_height, int(height * scale))
            resized = image if (scaled_width, scaled_height) == (width, height) else image.resize((scaled_width, scaled_height), Image.LANCZOS)

            for quality in quality_steps:
                candidate = path.with_name(f'{path.stem}-optimized-{scaled_width}x{scaled_height}-q{quality}.jpg')
                resized.save(candidate, format='JPEG', quality=quality, optimize=True, progressive=True)
                size = candidate.stat().st_size
                best_candidate = candidate
                if size <= max_bytes:
                    try:
                        path.unlink(missing_ok=True)
                    except Exception:
                        pass
                    return str(candidate)

        if best_candidate is not None:
            try:
                path.unlink(missing_ok=True)
            except Exception:
                pass
            return str(best_candidate)

    return png_path


def _convert_pptx_with_soffice(input_pptx: str, output_dir: str, target_format: str, soffice_path: str = '') -> str:
    soffice = resolve_soffice_path(soffice_path)
    output_root = Path(output_dir)
    output_root.mkdir(parents=True, exist_ok=True)
    profile_root = output_root / '_soffice_profiles'
    profile_root.mkdir(parents=True, exist_ok=True)
    profile_dir = profile_root / uuid4().hex
    profile_dir.mkdir(parents=True, exist_ok=True)
    profile_uri = profile_dir.resolve().as_uri()
    cmd = [
        soffice,
        '--headless',
        f'-env:UserInstallation={profile_uri}',
        '--convert-to',
        target_format,
        '--outdir',
        str(output_root),
        input_pptx,
    ]
    result = subprocess.run(cmd, capture_output=True, text=True, check=False)
    stdout = (result.stdout or '').strip()
    stderr = (result.stderr or '').strip()
    output_path = output_root / f'{Path(input_pptx).stem}.{target_format}'
    try:
        shutil.rmtree(profile_dir, ignore_errors=True)
    except Exception:
        pass
    if result.returncode != 0:
        details = stderr or stdout or 'tanpa output error dari LibreOffice'
        raise RuntimeError(
            f'Konversi {target_format.upper()} gagal. returncode={result.returncode}. '
            f'input={Path(input_pptx).name}. detail={details}'
        )
    if not output_path.exists():
        details = stderr or stdout or f'LibreOffice tidak menghasilkan file {target_format.upper()}'
        raise RuntimeError(
            f'File {target_format.upper()} hasil konversi tidak ditemukan. '
            f'input={Path(input_pptx).name}. detail={details}'
        )
    return str(output_path)


def build_pdf_safe_template(
    input_pptx: str,
    output_pptx: str,
    placeholders: list[str] | None = None,
    working_dir: str = '',
    soffice_path: str = '',
    cleanup_working_dir: bool = True,
) -> dict:
    source_path = Path(input_pptx)
    target_path = Path(output_pptx)
    placeholder_tokens = expand_placeholder_tokens(list(placeholders or []))
    if not placeholder_tokens:
        raise RuntimeError('Daftar placeholder template kosong. Tidak dapat menyiapkan mode aman PDF.')

    work_root = Path(working_dir) if working_dir else (target_path.parent / f'_{target_path.stem}_pdf_safe_build')
    shutil.rmtree(work_root, ignore_errors=True)
    work_root.mkdir(parents=True, exist_ok=True)

    render_input = work_root / f'{source_path.stem}-blank-placeholders.pptx'
    replace_placeholders(
        str(source_path),
        str(render_input),
        {token: '' for token in placeholder_tokens},
    )
    rendered_png = Path(_convert_pptx_with_soffice(str(render_input), str(work_root), 'png', soffice_path))
    optimized_background = Path(_compress_png_for_pdf_target(str(rendered_png)))

    src_prs = Presentation(str(source_path))
    placeholder_refs = _validate_placeholder_shapes(src_prs, placeholder_tokens)

    output_prs = Presentation()
    output_prs.slide_width = src_prs.slide_width
    output_prs.slide_height = src_prs.slide_height
    output_slide = output_prs.slides.add_slide(output_prs.slide_layouts[6])
    output_slide.shapes.add_picture(
        str(optimized_background),
        0,
        0,
        width=output_prs.slide_width,
        height=output_prs.slide_height,
    )

    copied_shapes = []
    for slide_index, shape_index in placeholder_refs:
        src_shape = src_prs.slides[slide_index].shapes[shape_index]
        textbox = output_slide.shapes.add_textbox(src_shape.left, src_shape.top, src_shape.width, src_shape.height)
        _copy_textbox_style(src_shape, textbox)
        copied_shapes.append({
            'name': getattr(src_shape, 'name', f'shape-{shape_index}'),
            'text': src_shape.text,
            'slide_index': slide_index,
            'shape_index': shape_index,
        })

    if len(output_prs.slides) > 1:
        rel_id = output_prs.slides._sldIdLst[0].rId
        output_prs.part.drop_rel(rel_id)
        del output_prs.slides._sldIdLst[0]

    target_path.parent.mkdir(parents=True, exist_ok=True)
    output_prs.save(str(target_path))

    background_png_path = str(optimized_background)
    working_dir_path = str(work_root)
    if cleanup_working_dir:
        try:
            shutil.rmtree(work_root, ignore_errors=True)
        except Exception:
            pass

    return {
        'output_pptx': str(target_path),
        'background_png': background_png_path,
        'placeholder_count': len(copied_shapes),
        'placeholders': copied_shapes,
        'working_dir': working_dir_path,
        'working_dir_cleaned': cleanup_working_dir,
    }


def convert_document_with_soffice(input_path: str, output_dir: str, target_format: str, soffice_path: str = '') -> str:
    return _convert_pptx_with_soffice(input_path, output_dir, target_format, soffice_path)


def convert_pptx_to_pdf_with_soffice(input_pptx: str, output_dir: str, soffice_path: str = '') -> str:
    return convert_document_with_soffice(input_pptx, output_dir, 'pdf', soffice_path)


def convert_pptx_to_pdf(input_pptx: str, output_dir: str) -> str:
    return convert_pptx_to_pdf_with_soffice(input_pptx, output_dir, current_app.config.get('SOFFICE_PATH', ''))
